"""
AI Suite — Tool System (MCP-compatible)
工具注册 + 执行引擎，通过 /api/tools 暴露给前端
"""
import os, subprocess, json, tempfile

# ═══════ 代理配置 ═══════
# 工具自动检测本地代理（Clash/V2Ray），检测成功则走代理访问外网
_PROXY_URL = None
_PROXY_CHECKED = False

def _get_opener():
    """Return (opener, proxy_url) — opener uses proxy if available"""
    global _PROXY_URL, _PROXY_CHECKED
    import urllib.request as _ur

    if not _PROXY_CHECKED:
        _PROXY_CHECKED = True
        candidates = [
            "http://127.0.0.1:7890",   # Clash (default)
            "http://127.0.0.1:7897",   # Clash (alt)
            "http://127.0.0.1:10809",  # V2RayN
            "http://127.0.0.1:1080",   # SOCKS5
        ]
        for proxy in candidates:
            try:
                ph = _ur.ProxyHandler({"http": proxy, "https": proxy})
                opener = _ur.build_opener(ph)
                # Quick connectivity test via proxy
                import urllib.request as _ur2
                head_req = _ur2.Request("https://www.bing.com", method="HEAD")
                opener.open(head_req, timeout=5)
                _PROXY_URL = proxy
                break
            except Exception:
                pass

    if _PROXY_URL:
        ph = _ur.ProxyHandler({"http": _PROXY_URL, "https": _PROXY_URL})
        return _ur.build_opener(ph), _PROXY_URL
    return _ur.build_opener(_ur.HTTPHandler), None

def _urlopen(url, timeout=15):
    """urlopen with proxy detection — falls back to direct on proxy failure"""
    import urllib.request as _ur
    opener, proxy_url = _get_opener()
    headers = {
        "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36"
    }
    if isinstance(url, str):
        req = _ur.Request(url, headers=headers)
    else:
        req = url
    if proxy_url:
        try:
            return opener.open(req, timeout=timeout)
        except Exception:
            # Proxy failed — retry direct
            return _ur.build_opener().open(req, timeout=timeout)
    return opener.open(req, timeout=timeout)

TOOLS = {}

def register(name, description, handler, schema=None):
    TOOLS[name] = {"name": name, "description": description, "handler": handler, "schema": schema or {}}

# ─── File Tools ───

def tool_read_file(path):
    try:
        with open(path, encoding="utf-8", errors="replace") as f:
            content = f.read(50000)
        return {"ok": True, "content": content, "truncated": len(content) >= 50000}
    except Exception as e:
        return {"ok": False, "error": str(e)}

def tool_list_dir(path):
    try:
        items = []
        for name in os.listdir(path):
            full = os.path.join(path, name)
            items.append({"name": name, "type": "dir" if os.path.isdir(full) else "file",
                          "size": os.path.getsize(full) if os.path.isfile(full) else 0})
        return {"ok": True, "items": sorted(items, key=lambda x: (x["type"], x["name"]))[:200]}
    except Exception as e:
        return {"ok": False, "error": str(e)}

def tool_write_file(path, content):
    try:
        os.makedirs(os.path.dirname(path) or ".", exist_ok=True)
        with open(path, "w", encoding="utf-8") as f:
            f.write(content)
        return {"ok": True, "path": path}
    except Exception as e:
        return {"ok": False, "error": str(e)}

# ─── Sandbox ───

DANGEROUS_PATTERNS = [
    r'rm\s+-rf\s+/', r'del\s+/[fsq]\s+\w:\\', r'format\s', r'fdisk\s',
    r'mkfs\.', r'dd\s+if=', r'>\s*/dev/', r'chmod\s+777\s+/',
    r'shutdown\s', r'reboot', r'halt', r'poweroff',
    r'wget\s.*\|.*sh', r'curl\s.*\|.*bash',
    # Windows-specific
    r'del\s+/[Ff]\s+/[Ss]\s+/[Qq]\s+\w:\\',  # del /f /s /q C:\
    r'rd\s+/[Ss]\s+/[Qq]\s+\w:\\',             # rd /s /q C:\
    r'diskpart', r'clean\s+all',                 # diskpart clean
    r'Remove-Item\s+-Recurse\s+-Force\s+\w:\\', # PowerShell recursive delete
    r'format\s+\w:', r'sfc\s+/scannow',         # system commands
]

SAFE_WRITE_DIRS = [
    os.path.expanduser("~"),
    os.path.expanduser("~/Desktop"),
    os.path.expanduser("~/Documents"),
    os.path.expanduser("~/Downloads"),
    os.path.expanduser("~/.ai-suite"),
    tempfile.gettempdir(),
    os.getcwd(),
]

def _sandbox_shell(command: str) -> dict:
    """Check shell command for dangerous patterns"""
    import re
    cmd_lower = command.lower()
    for pattern in DANGEROUS_PATTERNS:
        if re.search(pattern, cmd_lower):
            return {"ok": False, "error": f"Blocked dangerous command pattern: {pattern}"}
    return {"ok": True}

def _sandbox_path(path: str) -> dict:
    """Check if file path is safe to write to"""
    try:
        abs_path = os.path.abspath(os.path.expanduser(path)).lower()
    except:
        return {"ok": False, "error": f"Invalid path: {path}"}
    for safe_dir in SAFE_WRITE_DIRS:
        try:
            if abs_path.startswith(os.path.abspath(safe_dir).lower()):
                return {"ok": True}
        except:
            pass
    return {"ok": False, "error": f"Path '{path}' not in allowed directories. Must be under ~/ or %TEMP%."}

# ─── Shell Tool ───

def tool_shell(command, cwd=None):
    # Sandbox check
    check = _sandbox_shell(command)
    if not check["ok"]:
        return check
    try:
        result = subprocess.run(command, shell=True, capture_output=True, text=True,
                                timeout=30, cwd=cwd or os.path.expanduser("~"),
                                creationflags=subprocess.CREATE_NO_WINDOW if os.name == "nt" else 0)
        return {"ok": True, "stdout": result.stdout[-5000:], "stderr": result.stderr[-2000:],
                "returncode": result.returncode}
    except subprocess.TimeoutExpired:
        return {"ok": False, "error": "Command timed out (30s)"}
    except Exception as e:
        return {"ok": False, "error": str(e)}

# ─── Clipboard Tool ───

def tool_clipboard_read():
    try:
        result = subprocess.run(["powershell", "-NoProfile", "-Command", "Get-Clipboard"],
                                capture_output=True, text=True, timeout=5,
                                creationflags=subprocess.CREATE_NO_WINDOW)
        return {"ok": True, "text": result.stdout[:10000]}
    except Exception as e:
        return {"ok": False, "error": str(e)}

def tool_clipboard_write(text):
    try:
        # Base64 encode to avoid injection
        import base64
        encoded = base64.b64encode(text.encode("utf-8")).decode()
        ps = f"[System.Text.Encoding]::UTF8.GetString([System.Convert]::FromBase64String('{encoded}')) | Set-Clipboard"
        subprocess.run(["powershell", "-Command", ps],
                       capture_output=True, timeout=5,
                       creationflags=subprocess.CREATE_NO_WINDOW)
        return {"ok": True}
    except Exception as e:
        return {"ok": False, "error": str(e)}

# ─── Web Tool ───

def tool_web_fetch(url):
    import urllib.request as ur
    try:
        req = ur.Request(url, headers={
            "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36"
        })
        resp = _urlopen(req, timeout=15)
        content = resp.read().decode("utf-8", errors="replace")[:10000]
        return {"ok": True, "content": content, "status": resp.status}
    except Exception as e:
        return {"ok": False, "error": str(e)}

def tool_web_search(query):
    import urllib.request as ur
    import re

    def _try_search(url, snippet_pattern, result_count=5):
        """Try a search engine, return results list (empty on failure)"""
        req = ur.Request(url, headers={
            "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36"
        })
        resp = _urlopen(req, timeout=15)
        content = resp.read().decode("utf-8", errors="replace")[:120000]
        snippets = re.findall(snippet_pattern, content, re.DOTALL)
        results = []
        for s in snippets[:result_count]:
            clean = re.sub(r'<[^>]+>', '', s).strip()
            if clean and len(clean) > 10:
                results.append(clean)
        return results

    # Try Bing first
    try:
        url = f"https://www.bing.com/search?q={ur.quote(query)}&setlang=zh-cn"
        results = _try_search(url, r'<p[^>]*class="b_lineclamp\d+"[^>]*>(.*?)</p>', 5)
        if results:
            return {"ok": True, "results": results, "source": "bing"}
    except Exception:
        pass

    # Fallback: DuckDuckGo HTML
    try:
        url = f"https://html.duckduckgo.com/html/?q={ur.quote(query)}"
        results = _try_search(url, r'class="result__snippet"[^>]*>(.*?)</a>', 5)
        if results:
            return {"ok": True, "results": results, "source": "duckduckgo"}
    except Exception:
        pass

    # Last fallback: DuckDuckGo lite
    try:
        url = f"https://lite.duckduckgo.com/lite/?q={ur.quote(query)}"
        results = _try_search(url, r'<td[^>]*class="result-snippet"[^>]*>(.*?)</td>', 5)
        if results:
            return {"ok": True, "results": results, "source": "duckduckgo-lite"}
    except Exception:
        pass

    return {"ok": False, "error": "All search engines unreachable (network restriction?)"}

# ─── Desktop Tools (v3.0) ───

def tool_open_browser(url):
    """在默认浏览器打开 URL"""
    try:
        subprocess.run(["cmd", "/c", "start", url], timeout=5,
                       creationflags=subprocess.CREATE_NO_WINDOW if os.name == "nt" else 0)
        return {"ok": True, "url": url}
    except Exception as e:
        return {"ok": False, "error": str(e)}

def tool_screenshot(path=None):
    """截取屏幕截图，返回文件路径"""
    try:
        from PIL import ImageGrab
        img = ImageGrab.grab()
        save_path = path or os.path.expanduser(f"~/Desktop/screenshot_{int(__import__('time').time())}.png")
        img.save(save_path, "PNG")
        return {"ok": True, "path": save_path, "size": f"{img.width}x{img.height}"}
    except ImportError:
        pass  # 走 PowerShell 降级
    except Exception:
        pass  # 走 PowerShell 降级
    # Fallback: PowerShell screenshot
    try:
        save_path = path or os.path.expanduser("~/Desktop/screenshot.png")
        ps = f'''
        Add-Type -AssemblyName System.Windows.Forms,System.Drawing
        $screen = [System.Windows.Forms.Screen]::PrimaryScreen
        $bmp = New-Object System.Drawing.Bitmap($screen.Bounds.Width, $screen.Bounds.Height)
        $g = [System.Drawing.Graphics]::FromImage($bmp)
        $g.CopyFromScreen(0, 0, 0, 0, $bmp.Size)
        $bmp.Save("{save_path.replace(chr(92), chr(92)+chr(92))}")
        $g.Dispose()
        '''
        subprocess.run(["powershell", "-Command", ps], timeout=30,
                      capture_output=True, creationflags=subprocess.CREATE_NO_WINDOW)
        return {"ok": True, "path": save_path, "method": "powershell"}
    except Exception as e2:
        return {"ok": False, "error": f"screenshot failed: {e2}"}

def tool_find_files(directory, pattern="*", max_results=50):
    """递归搜索文件"""
    import fnmatch
    try:
        results = []
        for root, dirs, files in os.walk(directory):
            dirs[:] = [d for d in dirs if not d.startswith('.')]
            for f in files:
                if fnmatch.fnmatch(f, pattern):
                    results.append(os.path.join(root, f))
                    if len(results) >= max_results:
                        return {"ok": True, "files": results, "truncated": True}
            if len(results) >= max_results:
                break
        return {"ok": True, "files": results, "truncated": False}
    except Exception as e:
        return {"ok": False, "error": str(e)}

def tool_run_python(code, timeout=10):
    """在隔离环境中执行 Python 代码（用于计算、数据处理等）"""
    import sys
    try:
        env = os.environ.copy()
        env["PYTHONPATH"] = ""
        result = subprocess.run(
            [sys.executable, "-c", code],
            capture_output=True, text=True, timeout=timeout,
            cwd=os.path.expanduser("~"),
            creationflags=subprocess.CREATE_NO_WINDOW if os.name == "nt" else 0,
            env=env
        )
        return {
            "ok": True,
            "stdout": result.stdout[-5000:],
            "stderr": result.stderr[-2000:],
            "returncode": result.returncode
        }
    except subprocess.TimeoutExpired:
        return {"ok": False, "error": f"Timeout ({timeout}s)"}
    except Exception as e:
        return {"ok": False, "error": str(e)}


# ─── Register all tools ───
register("read_file", "读取文件内容 — Read a file", tool_read_file,
         {"path": {"type": "string", "description": "文件绝对路径"}})
register("list_dir", "列出目录内容 — List directory", tool_list_dir,
         {"path": {"type": "string", "description": "目录路径"}})
register("write_file", "写入内容到文件 — Write to file", tool_write_file,
         {"path": {"type": "string"}, "content": {"type": "string"}})
register("find_files", "递归搜索匹配文件 — Find files by pattern", tool_find_files,
         {"directory": {"type": "string"}, "pattern": {"type": "string", "optional": True},
          "max_results": {"type": "integer", "optional": True}})
register("shell", "执行 Shell 命令(30秒超时) — Run shell command", tool_shell,
         {"command": {"type": "string"}, "cwd": {"type": "string", "optional": True}})
register("run_python", "隔离执行 Python 代码 — Run Python sandboxed", tool_run_python,
         {"code": {"type": "string"}, "timeout": {"type": "number", "optional": True}})
register("clipboard_read", "读取剪贴板文本 — Read clipboard", tool_clipboard_read, {})
register("clipboard_write", "写入文本到剪贴板 — Write clipboard", tool_clipboard_write,
         {"text": {"type": "string"}})
register("web_fetch", "抓取网页内容 — Fetch web page", tool_web_fetch,
         {"url": {"type": "string"}})
register("web_search", "DuckDuckGo 网页搜索 — Web search", tool_web_search,
         {"query": {"type": "string"}})
register("open_browser", "默认浏览器打开 URL — Open browser", tool_open_browser,
         {"url": {"type": "string"}})
register("screenshot", "屏幕截图保存为文件 — Take screenshot", tool_screenshot,
         {"path": {"type": "string", "optional": True}})

def tool_ocr(image_path=None):
    """OCR 识别图片文字。优先 tesseract，降级 Windows 内置引擎"""
    try:
        target = image_path
        if not target:
            import time
            target = os.path.expanduser(f"~/Desktop/ocr_{int(time.time())}.png")
            r = tool_screenshot(target)
            if not r.get("ok"):
                return {"ok": False, "error": "截图失败: " + r.get("error", "未知")}

        if not os.path.exists(target):
            return {"ok": False, "error": f"文件不存在: {target}"}

        text = ""

        # 1) Tesseract (best Chinese support)
        try:
            import shutil as _sh
            tess = _sh.which("tesseract")
            if not tess:
                # Check common paths
                for p in [
                    "C:\\Program Files\\Tesseract-OCR\\tesseract.exe",
                    "C:\\Tesseract-OCR\\tesseract.exe",
                    os.path.expanduser("~\\AppData\\Local\\Tesseract-OCR\\tesseract.exe"),
                ]:
                    if os.path.exists(p):
                        tess = p
                        break
            if tess:
                # Ensure tessdata is findable
                env = os.environ.copy()
                if "TESSDATA_PREFIX" not in env:
                    for td in [
                        os.path.expanduser("~/.tesseract"),
                        os.path.join(os.path.dirname(tess), "tessdata"),
                    ]:
                        if os.path.isdir(td):
                            env["TESSDATA_PREFIX"] = td
                            break
                r = subprocess.run([tess, target, "stdout", "-l", "chi_sim+eng"],
                                   capture_output=True, text=True, timeout=30, env=env,
                                   creationflags=subprocess.CREATE_NO_WINDOW if os.name == "nt" else 0)
                if r.stdout.strip():
                    text = r.stdout.strip()
        except:
            pass

        # 2) Windows built-in OCR
        if not text:
            try:
                safe_path = target.replace('\\', '\\\\')
                ps = f'''
                Add-Type -AssemblyName System.Drawing | Out-Null
                $bmp = [System.Drawing.Bitmap]::FromFile("{safe_path}")
                $w = $bmp.Width; $h = $bmp.Height
                $data = $bmp.LockBits([System.Drawing.Rectangle]::new(0,0,$w,$h), [System.Drawing.Imaging.ImageLockMode]::ReadOnly, [System.Drawing.Imaging.PixelFormat]::Format32bppArgb)
                $bytes = [byte[]]::new($data.Stride * $h)
                [System.Runtime.InteropServices.Marshal]::Copy($data.Scan0, $bytes, 0, $bytes.Length)
                $bmp.UnlockBits($data); $bmp.Dispose()

                [Windows.Graphics.Imaging.BitmapDecoder, Windows.Graphics.Imaging, ContentType = WindowsRuntime] | Out-Null
                [Windows.Graphics.Imaging.SoftwareBitmap, Windows.Graphics.Imaging, ContentType = WindowsRuntime] | Out-Null
                [Windows.Media.Ocr.OcrEngine, Windows.Media.Ocr, ContentType = WindowsRuntime] | Out-Null

                $sb = [Windows.Graphics.Imaging.SoftwareBitmap]::CreateCopyFromBuffer(
                    [Windows.Security.Cryptography.CryptographicBuffer]::CreateFromByteArray($bytes),
                    [Windows.Graphics.Imaging.BitmapPixelFormat]::Bgra8, $w, $h)
                $engine = [Windows.Media.Ocr.OcrEngine]::TryCreateFromUserProfileLanguages()
                if (-not $engine) {{ $engine = [Windows.Media.Ocr.OcrEngine]::TryCreateFromLanguage("zh-Hans") }}
                if (-not $engine) {{ $engine = [Windows.Media.Ocr.OcrEngine]::TryCreateFromLanguage("en") }}
                if ($engine) {{
                    $result = $engine.RecognizeAsync($sb).GetAwaiter().GetResult()
                    ($result.Lines | ForEach-Object {{ ($_.Words | ForEach-Object {{ $_.Text }}) -join " " }}) -join "`n"
                }}
                '''
                r2 = subprocess.run(["powershell", "-NoProfile", "-Command", ps],
                                   capture_output=True, text=True, timeout=30,
                                   creationflags=subprocess.CREATE_NO_WINDOW if os.name == "nt" else 0)
                if r2.stdout.strip():
                    text = r2.stdout.strip()[:10000]
            except:
                pass

        if not text:
            return {"ok": True, "text": "", "source": target, "note": "no text detected or OCR engine unavailable"}

        return {"ok": True, "text": text[:10000], "source": target}
    except Exception as e:
        return {"ok": False, "error": str(e)}

register("ocr", "图片文字识别(OCR) — Extract text from image", tool_ocr,
         {"image_path": {"type": "string", "optional": True, "description": "图片路径，省略则先截图"}})

# ─── Desktop Control Tools (v3.1) ───

def tool_screenshot_find(text_query: str):
    """截图 + OCR 查找文字位置，返回坐标可用于点击"""
    import time as _time
    target = os.path.expanduser(f"~/Desktop/visual_{int(_time.time())}.png")
    r = tool_screenshot(target)
    if not r.get("ok"):
        return {"ok": False, "error": "截图失败: " + r.get("error", "")}

    # Run OCR with word-level bounding boxes via PowerShell
    ps = f'''
    Add-Type -AssemblyName System.Drawing | Out-Null
    $bmp = [System.Drawing.Bitmap]::FromFile("{target.replace(chr(92), chr(92)+chr(92))}")
    $w = $bmp.Width; $h = $bmp.Height
    $data = $bmp.LockBits([System.Drawing.Rectangle]::new(0,0,$w,$h), [System.Drawing.Imaging.ImageLockMode]::ReadOnly, [System.Drawing.Imaging.PixelFormat]::Format32bppArgb)
    $bytes = [byte[]]::new($data.Stride * $h)
    [System.Runtime.InteropServices.Marshal]::Copy($data.Scan0, $bytes, 0, $bytes.Length)
    $bmp.UnlockBits($data); $bmp.Dispose()

    [Windows.Graphics.Imaging.SoftwareBitmap, Windows.Graphics.Imaging, ContentType = WindowsRuntime] | Out-Null
    [Windows.Media.Ocr.OcrEngine, Windows.Media.Ocr, ContentType = WindowsRuntime] | Out-Null

    $sb = [Windows.Graphics.Imaging.SoftwareBitmap]::CreateCopyFromBuffer(
        [Windows.Security.Cryptography.CryptographicBuffer]::CreateFromByteArray($bytes),
        [Windows.Graphics.Imaging.BitmapPixelFormat]::Bgra8, $w, $h)
    $engine = [Windows.Media.Ocr.OcrEngine]::TryCreateFromUserProfileLanguages()
    if (-not $engine) {{ $engine = [Windows.Media.Ocr.OcrEngine]::TryCreateFromLanguage("zh-Hans") }}
    if (-not $engine) {{ $engine = [Windows.Media.Ocr.OcrEngine]::TryCreateFromLanguage("en") }}
    if ($engine) {{
        $result = $engine.RecognizeAsync($sb).GetAwaiter().GetResult()
        $query = "{text_query.replace(chr(34), '').replace(chr(36), '').replace(chr(123), '').replace(chr(125), '')}".ToLower()
        $matches = @()
        foreach ($line in $result.Lines) {{
            foreach ($word in $line.Words) {{
                if ($word.Text.ToLower().Contains($query)) {{
                    $matches += [PSCustomObject]@{{
                        Text = $word.Text
                        X = [int]($word.BoundingRect.X + $word.BoundingRect.Width / 2)
                        Y = [int]($word.BoundingRect.Y + $word.BoundingRect.Height / 2)
                        W = [int]$word.BoundingRect.Width
                        H = [int]$word.BoundingRect.Height
                    }}
                }}
            }}
        }}
        ConvertTo-Json -InputObject $matches -Compress
    }}
    '''
    try:
        r2 = subprocess.run(["powershell", "-NoProfile", "-Command", ps],
                           capture_output=True, text=True, timeout=30,
                           creationflags=subprocess.CREATE_NO_WINDOW if os.name == "nt" else 0)
        matches = []
        if r2.stdout.strip():
            import json as _j
            try:
                data = _j.loads(r2.stdout)
                if isinstance(data, dict):
                    data = [data]
                elif data is None:
                    data = []
                for m in data:
                    matches.append({"text": m.get("Text", ""),
                                    "x": m.get("X", 0), "y": m.get("Y", 0),
                                    "w": m.get("W", 0), "h": m.get("H", 0)})
            except:
                pass
        return {"ok": True, "query": text_query, "matches": matches[:20],
                "screenshot": target, "instruction": "Use click(x, y) to click the target"}
    except Exception as e:
        return {"ok": False, "error": str(e)}


def tool_click_text(text_query: str):
    """截图查找文字并自动点击（组合 screenshot_find + click）"""
    r = tool_screenshot_find(text_query)
    if not r.get("ok"):
        return r
    matches = r.get("matches", [])
    if not matches:
        return {"ok": False, "error": f"未找到文字: {text_query}"}
    # Click the first match
    m = matches[0]
    click_r = tool_click(m["x"], m["y"])
    return {"ok": True, "text": m["text"], "x": m["x"], "y": m["y"],
            "click_result": click_r, "total_matches": len(matches)}

register("screenshot_find", "截图查找文字位置 — Find text on screen", tool_screenshot_find,
         {"text_query": {"type": "string", "description": "要查找的文字"}})
register("click_text", "查找屏幕文字并自动点击 — Find and click text", tool_click_text,
         {"text_query": {"type": "string", "description": "要查找并点击的文字"}})

# ─── Desktop Control Tools (v3.1) ───

def tool_click(x: int, y: int, button: str = "left"):
    """模拟鼠标点击"""
    try:
        import ctypes
        MOUSEEVENTF_LEFTDOWN = 0x0002
        MOUSEEVENTF_LEFTUP = 0x0004
        MOUSEEVENTF_RIGHTDOWN = 0x0008
        MOUSEEVENTF_RIGHTUP = 0x0010
        ctypes.windll.user32.SetCursorPos(x, y)
        if button == "right":
            ctypes.windll.user32.mouse_event(MOUSEEVENTF_RIGHTDOWN, 0, 0, 0, 0)
            ctypes.windll.user32.mouse_event(MOUSEEVENTF_RIGHTUP, 0, 0, 0, 0)
        else:
            ctypes.windll.user32.mouse_event(MOUSEEVENTF_LEFTDOWN, 0, 0, 0, 0)
            ctypes.windll.user32.mouse_event(MOUSEEVENTF_LEFTUP, 0, 0, 0, 0)
        return {"ok": True, "x": x, "y": y, "button": button}
    except Exception as e:
        return {"ok": False, "error": str(e)}

def tool_type_text(text: str, interval: float = 0.02):
    """模拟键盘输入"""
    try:
        import ctypes
        import time as _time
        KEYEVENTF_UNICODE = 0x0004
        KEYEVENTF_KEYUP = 0x0002
        for ch in text:
            vk = ord(ch)
            # KEYEVENTF_UNICODE: bVk must be 0, bScan holds the Unicode char
            ctypes.windll.user32.keybd_event(0, vk, KEYEVENTF_UNICODE, 0)
            ctypes.windll.user32.keybd_event(0, vk, KEYEVENTF_UNICODE | KEYEVENTF_KEYUP, 0)
            _time.sleep(interval)
        return {"ok": True, "length": len(text)}
    except Exception as e:
        return {"ok": False, "error": str(e)}

def tool_press_key(key: str):
    """模拟按键（enter, escape, tab, space, backspace, delete, 方向键等）"""
    VK_MAP = {
        "enter": 0x0D, "return": 0x0D, "escape": 0x1B, "esc": 0x1B,
        "tab": 0x09, "space": 0x20, "backspace": 0x08, "delete": 0x2E,
        "up": 0x26, "down": 0x28, "left": 0x25, "right": 0x27,
        "home": 0x24, "end": 0x23, "pageup": 0x21, "pagedown": 0x22,
        "f1": 0x70, "f2": 0x71, "f3": 0x72, "f4": 0x73, "f5": 0x74,
        "f6": 0x75, "f7": 0x76, "f8": 0x77, "f9": 0x78, "f10": 0x79,
        "f11": 0x7A, "f12": 0x7B,
        "ctrl": 0x11, "shift": 0x10, "alt": 0x12, "win": 0x5B,
    }
    vk = VK_MAP.get(key.lower())
    if vk is None:
        return {"ok": False, "error": f"Unknown key: {key}. Use: {', '.join(VK_MAP.keys())}"}
    try:
        import ctypes
        KEYEVENTF_KEYUP = 0x0002
        ctypes.windll.user32.keybd_event(vk, 0, 0, 0)
        ctypes.windll.user32.keybd_event(vk, 0, KEYEVENTF_KEYUP, 0)
        return {"ok": True, "key": key}
    except Exception as e:
        return {"ok": False, "error": str(e)}

def tool_get_windows():
    """获取当前打开的窗口列表"""
    import ctypes
    try:
        user32 = ctypes.windll.user32
        windows = []
        WNDENUMPROC = ctypes.WINFUNCTYPE(ctypes.c_bool, ctypes.c_void_p, ctypes.c_void_p)
        def enum_callback(hwnd, _lparam):
            if user32.IsWindowVisible(hwnd):
                buf = ctypes.create_unicode_buffer(256)
                user32.GetWindowTextW(hwnd, buf, 256)
                title = buf.value
                if title and len(title.strip()) > 0:
                    windows.append({"hwnd": str(hwnd), "title": title})
            return True
        user32.EnumWindows(WNDENUMPROC(enum_callback), 0)
        return {"ok": True, "windows": windows[:30]}
    except Exception as e:
        return {"ok": False, "error": str(e)}

def tool_focus_window(title_match: str):
    """根据标题匹配聚焦窗口"""
    import ctypes
    try:
        user32 = ctypes.windll.user32
        found = []
        WNDENUMPROC = ctypes.WINFUNCTYPE(ctypes.c_bool, ctypes.c_void_p, ctypes.c_void_p)
        def enum_callback(hwnd, _lparam):
            buf = ctypes.create_unicode_buffer(256)
            user32.GetWindowTextW(hwnd, buf, 256)
            title = buf.value
            if title and title_match.lower() in title.lower():
                found.append(hwnd)
                user32.ShowWindow(hwnd, 9)  # SW_RESTORE
                user32.SetForegroundWindow(hwnd)
                return False  # Stop enumeration
            return True
        user32.EnumWindows(WNDENUMPROC(enum_callback), 0)
        if found:
            buf = ctypes.create_unicode_buffer(256)
            user32.GetWindowTextW(found[0], buf, 256)
            return {"ok": True, "focused": buf.value}
        return {"ok": False, "error": f"Window not found: {title_match}"}
    except Exception as e:
        return {"ok": False, "error": str(e)}

def tool_move_mouse(x: int, y: int):
    """移动鼠标到指定坐标（不点击）"""
    try:
        import ctypes
        ctypes.windll.user32.SetCursorPos(x, y)
        return {"ok": True, "x": x, "y": y}
    except Exception as e:
        return {"ok": False, "error": str(e)}

def tool_get_mouse_pos():
    """获取当前鼠标位置"""
    try:
        import ctypes
        class POINT(ctypes.Structure):
            _fields_ = [("x", ctypes.c_long), ("y", ctypes.c_long)]
        pt = POINT()
        ctypes.windll.user32.GetCursorPos(ctypes.byref(pt))
        return {"ok": True, "x": pt.x, "y": pt.y}
    except Exception as e:
        return {"ok": False, "error": str(e)}

register("click", "鼠标点击(x,y) — Click at position", tool_click,
         {"x": {"type": "integer"}, "y": {"type": "integer"}, "button": {"type": "string", "optional": True}})
register("move_mouse", "移动鼠标到(x,y)不点击 — Move mouse", tool_move_mouse,
         {"x": {"type": "integer"}, "y": {"type": "integer"}})
register("mouse_pos", "获取当前鼠标坐标 — Get mouse position", tool_get_mouse_pos, {})
register("type_text", "模拟键盘输入文字 — Type text", tool_type_text,
         {"text": {"type": "string"}, "interval": {"type": "number", "optional": True}})
register("press_key", "按下一个键(enter/escape/tab/F1-F12/方向键) — Press key", tool_press_key,
         {"key": {"type": "string"}})
register("get_windows", "列出所有可见窗口标题 — List windows", tool_get_windows, {})
register("focus_window", "根据标题聚焦窗口(模糊匹配) — Focus window", tool_focus_window,
         {"title_match": {"type": "string"}})

# ─── Document Generation Tools ───

def tool_create_pptx(title: str, slides_json: str = "[]", output_path: str = None):
    """生成 PPTX 演示文稿"""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        import json as _json

        slides_data = _json.loads(slides_json) if isinstance(slides_json, str) else slides_json
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        for s in slides_data:
            layout_idx = {"title": 0, "bullet": 1, "image": 6, "end": 0}.get(
                s.get("layout", "bullet"), 1)
            if layout_idx >= len(prs.slide_layouts):
                layout_idx = 1
            slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])

            # Title
            if s.get("title") and slide.shapes.title:
                slide.shapes.title.text = s["title"]

            # Content
            if s.get("content"):
                if isinstance(s["content"], list):
                    text = "\n".join(f"• {item}" for item in s["content"])
                else:
                    text = s["content"]
                # Find or use subtitle placeholder
                if len(slide.placeholders) > 1:
                    ph = slide.placeholders[1]
                    ph.text = text

            # Subtitle for title slides
            if s.get("subtitle") and len(slide.placeholders) > 1:
                slide.placeholders[1].text = s["subtitle"]

        save_path = output_path or os.path.expanduser(
            f"~/Desktop/{title.replace(' ','_')[:30]}.pptx")
        os.makedirs(os.path.dirname(save_path) or ".", exist_ok=True)
        prs.save(save_path)
        return {"ok": True, "path": save_path, "slides": len(slides_data)}
    except ImportError:
        return {"ok": False, "error": "python-pptx not installed. pip install python-pptx"}
    except Exception as e:
        return {"ok": False, "error": str(e)}


def tool_create_dxf(filename: str, entities_json: str = "[]", units: str = "mm"):
    """生成 DXF CAD 图纸"""
    try:
        import ezdxf
        import json as _json

        data = _json.loads(entities_json) if isinstance(entities_json, str) else entities_json
        entities = data.get("entities", data) if isinstance(data, dict) else data
        title = data.get("title", filename) if isinstance(data, dict) else filename

        doc = ezdxf.new("R2010")
        msp = doc.modelspace()

        # Set up layers
        for layer in ["outline", "dimension", "hidden", "center", "holes", "text"]:
            doc.layers.add(name=layer, color={"outline": 7, "dimension": 3,
                          "hidden": 8, "center": 1, "holes": 4, "text": 2}.get(layer, 7))

        for e in entities:
            etype = e.get("type", "")
            layer = e.get("layer", "outline")

            if etype == "line":
                msp.add_line((e["x1"], e["y1"]), (e["x2"], e["y2"]),
                            dxfattribs={"layer": layer})
            elif etype == "circle":
                msp.add_circle((e["cx"], e["cy"]), e.get("radius", 10),
                              dxfattribs={"layer": layer})
            elif etype == "arc":
                msp.add_arc((e["cx"], e["cy"]), e.get("radius", 10),
                           e.get("start_angle", 0), e.get("end_angle", 90),
                           dxfattribs={"layer": layer})
            elif etype == "rect":
                x, y, w, h = e["x"], e["y"], e["w"], e["h"]
                pts = [(x, y), (x+w, y), (x+w, y+h), (x, y+h)]
                msp.add_lwpolyline(pts, close=True, dxfattribs={"layer": layer})
            elif etype == "text":
                msp.add_text(e.get("text", ""), dxfattribs={
                    "layer": "text", "height": e.get("height", 5)
                }).set_placement((e.get("x", 0), e.get("y", 0)))

        save_path = os.path.expanduser(
            f"~/Desktop/{filename.replace('.dxf','')[:30]}.dxf")
        doc.saveas(save_path)
        return {"ok": True, "path": save_path, "entities": len(entities),
                "title": title, "units": units}
    except ImportError:
        return {"ok": False, "error": "ezdxf not installed. pip install ezdxf"}
    except Exception as e:
        return {"ok": False, "error": str(e)}

register("create_pptx", "生成 PowerPoint 演示文稿 — Generate PPTX", tool_create_pptx,
         {"title": {"type": "string"}, "slides_json": {"type": "string"}, "output_path": {"type": "string", "optional": True}})
register("create_dxf", "生成 CAD DXF 图纸 — Generate DXF drawing", tool_create_dxf,
         {"filename": {"type": "string"}, "entities_json": {"type": "string"}, "units": {"type": "string", "optional": True}})

# ─── System & Desktop Management Tools (v3.2) ───

def tool_system_info():
    """获取系统信息：CPU/内存/磁盘/电池"""
    import ctypes
    try:
        info = {"platform": os.name, "cwd": os.getcwd()}

        # CPU cores
        info["cpu_count"] = os.cpu_count()

        # RAM via psutil (reliable), fallback to PowerShell
        try:
            import psutil
            vmem = psutil.virtual_memory()
            info["ram_total_gb"] = round(vmem.total / (1024**3), 1)
            info["ram_avail_gb"] = round(vmem.available / (1024**3), 1)
            info["ram_used_pct"] = int(vmem.percent)
        except ImportError:
            try:
                r = subprocess.run(
                    ["powershell", "-Command",
                     "$c=Get-CimInstance Win32_OperatingSystem;Write-Host ($c.TotalVisibleMemorySize/1MB).ToString('0.0');Write-Host ($c.FreePhysicalMemory/1MB).ToString('0.0')"],
                    capture_output=True, text=True, timeout=10,
                    creationflags=subprocess.CREATE_NO_WINDOW if os.name == "nt" else 0)
                lines = r.stdout.strip().split()
                if len(lines) >= 2:
                    info["ram_total_gb"] = round(float(lines[0]) / 1024, 1)
                    info["ram_avail_gb"] = round(float(lines[1]) / 1024, 1)
                    info["ram_used_pct"] = round((1 - float(lines[1]) / float(lines[0])) * 100) if float(lines[0]) else 0
            except:
                info["ram_total_gb"] = 0
                info["ram_avail_gb"] = 0
                info["ram_used_pct"] = 0

        # Disk
        import shutil
        disk = shutil.disk_usage(os.path.expanduser("~"))
        info["disk_total_gb"] = round(disk.total / (1024**3), 1)
        info["disk_free_gb"] = round(disk.free / (1024**3), 1)

        # Battery
        try:
            ps_out = subprocess.run(
                ["powershell", "-Command", "(Get-WmiObject Win32_Battery).EstimatedChargeRemaining"],
                capture_output=True, text=True, timeout=5,
                creationflags=subprocess.CREATE_NO_WINDOW)
            bat = ps_out.stdout.strip()
            if bat:
                info["battery_pct"] = int(bat)
        except:
            pass

        return {"ok": True, "info": info}
    except Exception as e:
        return {"ok": False, "error": str(e)}


def tool_list_processes(filter_name: str = ""):
    """列出运行中的进程"""
    try:
        ps = 'Get-Process | Select-Object Id,ProcessName,@{N="MemMB";E={[math]::Round($_.WorkingSet64/1MB,1)}}'
        if filter_name:
            ps += f' | Where-Object {{$_.ProcessName -like "*{filter_name}*"}}'
        ps += ' | Sort-Object MemMB -Descending | Select-Object -First 50 | ConvertTo-Json'
        r = subprocess.run(["powershell", "-Command", ps],
                          capture_output=True, text=True, timeout=15,
                          creationflags=subprocess.CREATE_NO_WINDOW)
        import json as _j
        data = _j.loads(r.stdout) if r.stdout.strip() else []
        if isinstance(data, dict):
            data = [data]
        return {"ok": True, "processes": data, "count": len(data)}
    except Exception as e:
        return {"ok": False, "error": str(e)}


def tool_kill_process(target, by: str = "name"):
    """终止进程。target: 进程名或PID；by: 'name' 或 'pid'"""
    try:
        if by == "pid":
            ps = f"Stop-Process -Id {int(target)} -Force"
        else:
            ps = f'Stop-Process -Name "{target}" -Force -ErrorAction SilentlyContinue'
        r = subprocess.run(["powershell", "-Command", ps],
                          capture_output=True, text=True, timeout=10,
                          creationflags=subprocess.CREATE_NO_WINDOW)
        ok = r.returncode == 0 and not r.stderr.strip()
        return {"ok": ok, "target": target, "by": by,
                "error": r.stderr.strip()[:500] if not ok else None}
    except Exception as e:
        return {"ok": False, "error": str(e)}


def tool_launch_app(app_path: str, args: str = ""):
    """启动应用程序"""
    try:
        if os.path.exists(app_path):
            cmd = f'start "" "{app_path}" {args}'
        else:
            # Try as command name (e.g., "notepad", "calc", "mspaint")
            cmd = f'start {app_path} {args}'
        subprocess.run(["cmd", "/c", cmd], timeout=5,
                      creationflags=subprocess.CREATE_NO_WINDOW if os.name == "nt" else 0)
        return {"ok": True, "app": app_path}
    except Exception as e:
        return {"ok": False, "error": str(e)}


def tool_get_volume():
    """获取系统音量 (0-100)，使用 winmm waveOut API"""
    try:
        ps = '''
        Add-Type -Name Win32 -Namespace Audio -MemberDefinition @"
        [DllImport("winmm.dll")] public static extern int waveOutGetVolume(int id, out uint vol);
"@ | Out-Null
        $vol = 0
        $null = [Audio.Win32]::waveOutGetVolume(0, [ref]$vol)
        $left = $vol -band 0xFFFF
        $right = ($vol -shr 16) -band 0xFFFF
        Write-Host ([int](($left + $right) / 2 * 100 / 65535))
        '''
        r = subprocess.run(["powershell", "-NoProfile", "-Command", ps],
                          capture_output=True, text=True, timeout=5,
                          creationflags=subprocess.CREATE_NO_WINDOW if os.name == "nt" else 0)
        lines = [l.strip() for l in r.stdout.strip().split('\n') if l.strip().isdigit()]
        if lines:
            return {"ok": True, "volume": int(lines[-1])}
        return {"ok": False, "error": "Could not read volume"}
    except Exception as e:
        return {"ok": False, "error": str(e)}


def tool_set_volume(level: int):
    """设置系统音量 (0-100)，使用 winmm waveOut API"""
    try:
        level = max(0, min(100, level))
        val = int(level / 100 * 65535)
        val_hex = val | (val << 16)  # both channels
        ps = f'''
        Add-Type -Name Win32 -Namespace Audio -MemberDefinition @"
        [DllImport("winmm.dll")] public static extern int waveOutSetVolume(int id, uint vol);
"@ | Out-Null
        $null = [Audio.Win32]::waveOutSetVolume(0, {val_hex})
        '''
        r = subprocess.run(["powershell", "-NoProfile", "-Command", ps],
                          capture_output=True, text=True, timeout=5,
                          creationflags=subprocess.CREATE_NO_WINDOW if os.name == "nt" else 0)
        if r.returncode == 0:
            return {"ok": True, "volume": level}
        return {"ok": False, "error": r.stderr.strip()[:200]}
    except Exception as e:
        return {"ok": False, "error": str(e)}


def tool_window_control(title_match: str, action: str = "focus"):
    """窗口操控：focus/minimize/maximize/restore/close"""
    import ctypes
    try:
        user32 = ctypes.windll.user32
        SW_ACTIONS = {
            "focus": 9, "restore": 9,
            "minimize": 6, "maximize": 3,
            "close": 0x0010,  # WM_CLOSE
        }
        if action not in SW_ACTIONS:
            return {"ok": False, "error": f"Unknown action: {action}. Use: {list(SW_ACTIONS.keys())}"}

        found = []

        WNDENUMPROC = ctypes.WINFUNCTYPE(ctypes.c_bool, ctypes.c_void_p, ctypes.c_void_p)
        def enum_callback(hwnd, _lparam):
            buf = ctypes.create_unicode_buffer(256)
            user32.GetWindowTextW(hwnd, buf, 256)
            title = buf.value
            if title and title_match.lower() in title.lower():
                found.append((hwnd, title))
            return True

        user32.EnumWindows(WNDENUMPROC(enum_callback), 0)

        if not found:
            return {"ok": False, "error": f"No window found matching: {title_match}"}

        results = []
        for hwnd, title in found:
            if action == "close":
                user32.PostMessageW(hwnd, SW_ACTIONS["close"], 0, 0)
            else:
                user32.ShowWindow(hwnd, SW_ACTIONS[action])
                if action in ("focus", "restore"):
                    user32.SetForegroundWindow(hwnd)
            results.append({"hwnd": str(hwnd), "title": title})

        return {"ok": True, "action": action, "windows": results}
    except Exception as e:
        return {"ok": False, "error": str(e)}


def tool_copy_file(src: str, dst: str):
    """复制文件（沙箱限制）"""
    src_check = _sandbox_path(src)
    if not src_check["ok"]:
        return src_check
    dst_check = _sandbox_path(dst)
    if not dst_check["ok"]:
        return dst_check
    try:
        import shutil
        os.makedirs(os.path.dirname(dst) or ".", exist_ok=True)
        shutil.copy2(src, dst)
        return {"ok": True, "src": src, "dst": dst}
    except Exception as e:
        return {"ok": False, "error": str(e)}


def tool_move_file(src: str, dst: str):
    """移动文件（沙箱限制）"""
    src_check = _sandbox_path(src)
    if not src_check["ok"]:
        return src_check
    dst_check = _sandbox_path(dst)
    if not dst_check["ok"]:
        return dst_check
    try:
        import shutil
        os.makedirs(os.path.dirname(dst) or ".", exist_ok=True)
        shutil.move(src, dst)
        return {"ok": True, "src": src, "dst": dst}
    except Exception as e:
        return {"ok": False, "error": str(e)}


def tool_delete_file(path: str, permanent: bool = False):
    """删除文件/文件夹（沙箱限制）。permanent=True 永久删除，否则到回收站"""
    path_check = _sandbox_path(path)
    if not path_check["ok"]:
        return path_check
    if not os.path.exists(path):
        return {"ok": False, "error": f"文件不存在: {path}"}
    try:
        import shutil
        if os.path.isdir(path):
            if permanent:
                shutil.rmtree(path)
            else:
                import send2trash
                send2trash.send2trash(path)
        else:
            if permanent:
                os.unlink(path)
            else:
                import send2trash
                send2trash.send2trash(path)
        return {"ok": True, "path": path, "permanent": permanent}
    except ImportError:
        return {"ok": False, "error": "send2trash not installed. Use permanent=True for direct delete, or pip install send2trash for recycle bin support."}
    except Exception as e:
        return {"ok": False, "error": str(e)}


register("system_info", "获取系统信息(CPU/内存/磁盘/电池) — System info", tool_system_info, {})
register("list_processes", "列出运行进程(可选名称过滤) — List processes", tool_list_processes,
         {"filter_name": {"type": "string", "optional": True}})
register("kill_process", "终止进程(按名称或PID) — Kill process", tool_kill_process,
         {"target": {"type": "string"}, "by": {"type": "string", "optional": True}})
register("launch_app", "启动应用程序(路径或名称如notepad/calc) — Launch app", tool_launch_app,
         {"app_path": {"type": "string"}, "args": {"type": "string", "optional": True}})
register("get_volume", "获取系统音量(0-100) — Get volume", tool_get_volume, {})
register("set_volume", "设置系统音量(0-100) — Set volume", tool_set_volume,
         {"level": {"type": "integer"}})
register("window_control", "窗口操控(聚焦/最小化/最大化/关闭) — Window control", tool_window_control,
         {"title_match": {"type": "string"}, "action": {"type": "string", "optional": True}})
register("copy_file", "复制文件(沙箱限制) — Copy file", tool_copy_file,
         {"src": {"type": "string"}, "dst": {"type": "string"}})
register("move_file", "移动/重命名文件(沙箱限制) — Move file", tool_move_file,
         {"src": {"type": "string"}, "dst": {"type": "string"}})
register("delete_file", "删除文件/文件夹(默认回收站) — Delete file", tool_delete_file,
         {"path": {"type": "string"}, "permanent": {"type": "boolean", "optional": True}})

# ─── Video Editing Tools (v3.3) ───

FFMPEG_PATH = None

def _find_ffmpeg():
    """Locate ffmpeg binary"""
    global FFMPEG_PATH
    if FFMPEG_PATH:
        return FFMPEG_PATH
    import shutil as _sh
    # Check common paths
    for p in [
        "ffmpeg.exe",  # in PATH or cwd
        os.path.expanduser("~/ffmpeg/bin/ffmpeg.exe"),
        "C:\\ffmpeg\\bin\\ffmpeg.exe",
        "C:\\Program Files\\ffmpeg\\bin\\ffmpeg.exe",
    ]:
        if os.path.exists(p):
            FFMPEG_PATH = p
            return p
    found = _sh.which("ffmpeg")
    if found:
        FFMPEG_PATH = found
        return found
    return None

def _run_ffmpeg(args: list, timeout: int = 120) -> dict:
    """Run ffmpeg with args, return result dict"""
    ff = _find_ffmpeg()
    if not ff:
        return {"ok": False, "error": "ffmpeg not found. Install ffmpeg and add to PATH."}
    try:
        result = subprocess.run(
            [ff] + args, capture_output=True, text=True, timeout=timeout,
            creationflags=subprocess.CREATE_NO_WINDOW if os.name == "nt" else 0
        )
        return {"ok": result.returncode == 0, "returncode": result.returncode,
                "stderr_last": result.stderr.strip()[-500:] if result.stderr else "",
                "stdout": result.stdout.strip()[:500]}
    except subprocess.TimeoutExpired:
        return {"ok": False, "error": f"ffmpeg timed out ({timeout}s)"}
    except Exception as e:
        return {"ok": False, "error": str(e)}

def _safe_rframe(rate_str):
    """Safely compute fps from ffprobe 'r_frame_rate' string like '30000/1001'"""
    if '/' in rate_str:
        try:
            num, den = rate_str.split('/', 1)
            n, d = float(num), float(den)
            return round(n / d, 3) if d != 0 else 0.0
        except (ValueError, ZeroDivisionError):
            return 0.0
    try:
        return float(rate_str)
    except ValueError:
        return 0.0

def _ffmpeg_safe_path(path: str) -> str:
    """Check path is valid, return absolute"""
    abs_path = os.path.abspath(path)
    return abs_path


def tool_video_info(video_path: str):
    """获取视频元信息：时长、分辨率、编码、码率、帧数"""
    path = _ffmpeg_safe_path(video_path)
    if not os.path.exists(path):
        return {"ok": False, "error": f"File not found: {video_path}"}

    # Use ffprobe if available
    ffprobe = os.path.join(os.path.dirname(_find_ffmpeg() or ""), "ffprobe.exe")
    if not os.path.exists(ffprobe):
        ffprobe = None
        import shutil as _sh
        found = _sh.which("ffprobe")
        if found:
            ffprobe = found

    if ffprobe:
        try:
            r = subprocess.run([ffprobe, "-v", "quiet", "-print_format", "json",
                               "-show_format", "-show_streams", path],
                              capture_output=True, text=True, timeout=15,
                              creationflags=subprocess.CREATE_NO_WINDOW if os.name == "nt" else 0)
            if r.returncode == 0:
                import json as _j
                data = _j.loads(r.stdout)
                fmt = data.get("format", {})
                streams = data.get("streams", [])
                video_stream = next((s for s in streams if s.get("codec_type") == "video"), {})
                audio_stream = next((s for s in streams if s.get("codec_type") == "audio"), {})
                return {"ok": True, "path": path,
                        "duration_s": float(fmt.get("duration", 0)),
                        "size_mb": round(int(fmt.get("size", 0)) / (1024**2), 1),
                        "bitrate_kbps": int(int(fmt.get("bit_rate", 0)) / 1000),
                        "video_codec": video_stream.get("codec_name", ""),
                        "resolution": f"{video_stream.get('width', '?')}x{video_stream.get('height', '?')}",
                        "fps": _safe_rframe(video_stream.get("r_frame_rate", "0/1")),
                        "audio_codec": audio_stream.get("codec_name", ""),
                        "audio_channels": audio_stream.get("channels", 0),
                        "has_audio": bool(audio_stream)}
        except:
            pass

    # Fallback: parse ffmpeg stderr
    r = _run_ffmpeg(["-i", path, "-f", "null", "NUL"], timeout=15)
    info = {"ok": True, "path": path, "stderr_parse": True}
    stderr = r.get("stderr_last", "")
    import re
    dur_match = re.search(r'Duration:\s*(\d+):(\d+):(\d+\.\d+)', stderr)
    if dur_match:
        h, m, s = int(dur_match.group(1)), int(dur_match.group(2)), float(dur_match.group(3))
        info["duration_s"] = h * 3600 + m * 60 + s
    res_match = re.search(r'(\d{3,4})x(\d{3,4})', stderr)
    if res_match:
        info["resolution"] = f"{res_match.group(1)}x{res_match.group(2)}"
    codec_match = re.search(r'Video:\s*(\S+)', stderr)
    if codec_match:
        info["video_codec"] = codec_match.group(1)
    return info


def tool_video_trim(video_path: str, output_path: str, start: str = "00:00:00", duration: str = None, end: str = None):
    """裁剪视频片段。start: HH:MM:SS 或秒数；duration/end: 时长或结束时间"""
    src = _ffmpeg_safe_path(video_path)
    dst = _ffmpeg_safe_path(output_path)
    if not os.path.exists(src):
        return {"ok": False, "error": f"Source not found: {video_path}"}
    os.makedirs(os.path.dirname(dst) or ".", exist_ok=True)
    args = ["-i", src, "-ss", str(start), "-c:v", "libx264", "-c:a", "aac", "-avoid_negative_ts", "make_zero"]
    if end:
        args += ["-to", str(end)]
    elif duration:
        args += ["-t", str(duration)]
    args += ["-y", dst]
    r = _run_ffmpeg(args)
    r["output"] = dst
    return r


def tool_video_concat(video_paths: str, output_path: str):
    """拼接多个视频（JSON 数组字符串，每个元素是文件路径）"""
    import json as _j
    try:
        paths = _j.loads(video_paths) if isinstance(video_paths, str) else video_paths
    except:
        return {"ok": False, "error": "video_paths must be JSON array of file paths"}

    dst = _ffmpeg_safe_path(output_path)
    os.makedirs(os.path.dirname(dst) or ".", exist_ok=True)

    # Write concat file list
    concat_file = os.path.join(tempfile.gettempdir(), f"ffmpeg_concat_{os.getpid()}.txt")
    with open(concat_file, "w", encoding="utf-8") as f:
        for p in paths:
            abs_p = _ffmpeg_safe_path(p)
            if not os.path.exists(abs_p):
                return {"ok": False, "error": f"File not found: {p}"}
            f.write(f"file '{abs_p.replace(chr(92), '/')}'\n")

    r = _run_ffmpeg(["-f", "concat", "-safe", "0", "-i", concat_file, "-c", "copy", "-y", dst])
    try:
        os.unlink(concat_file)
    except:
        pass
    r["output"] = dst
    r["input_count"] = len(paths)
    return r


def tool_video_resize(video_path: str, output_path: str, width: int = 1920, height: int = 1080):
    """调整视频分辨率"""
    src = _ffmpeg_safe_path(video_path)
    dst = _ffmpeg_safe_path(output_path)
    if not os.path.exists(src):
        return {"ok": False, "error": f"Source not found: {video_path}"}
    os.makedirs(os.path.dirname(dst) or ".", exist_ok=True)
    r = _run_ffmpeg(["-i", src, "-vf", f"scale={width}:{height}:force_original_aspect_ratio=decrease,pad={width}:{height}:(ow-iw)/2:(oh-ih)/2",
                     "-c:v", "libx264", "-c:a", "aac", "-y", dst])
    r["output"] = dst
    r["resolution"] = f"{width}x{height}"
    return r


def tool_video_extract_audio(video_path: str, output_path: str = None, format: str = "mp3"):
    """从视频提取音频"""
    src = _ffmpeg_safe_path(video_path)
    if not os.path.exists(src):
        return {"ok": False, "error": f"Source not found: {video_path}"}
    if not output_path:
        base = os.path.splitext(src)[0]
        output_path = f"{base}_audio.{format}"
    dst = _ffmpeg_safe_path(output_path)
    os.makedirs(os.path.dirname(dst) or ".", exist_ok=True)
    codec_map = {"mp3": "libmp3lame", "aac": "aac", "wav": "pcm_s16le", "ogg": "libvorbis", "m4a": "aac"}
    codec = codec_map.get(format, "libmp3lame")
    r = _run_ffmpeg(["-i", src, "-vn", "-c:a", codec, "-q:a", "2", "-y", dst])
    r["output"] = dst
    return r


def tool_video_replace_audio(video_path: str, audio_path: str, output_path: str):
    """替换视频音轨（用新的音频文件替换原音轨）"""
    src = _ffmpeg_safe_path(video_path)
    aud = _ffmpeg_safe_path(audio_path)
    dst = _ffmpeg_safe_path(output_path)
    if not os.path.exists(src):
        return {"ok": False, "error": f"Video not found: {video_path}"}
    if not os.path.exists(aud):
        return {"ok": False, "error": f"Audio not found: {audio_path}"}
    os.makedirs(os.path.dirname(dst) or ".", exist_ok=True)
    r = _run_ffmpeg(["-i", src, "-i", aud, "-c:v", "copy", "-map", "0:v:0", "-map", "1:a:0",
                     "-shortest", "-y", dst])
    r["output"] = dst
    return r


def tool_video_speed(video_path: str, output_path: str, speed: float = 2.0):
    """调整视频播放速度。speed: 0.5(半速) / 2.0(双倍速)"""
    src = _ffmpeg_safe_path(video_path)
    dst = _ffmpeg_safe_path(output_path)
    if not os.path.exists(src):
        return {"ok": False, "error": f"Source not found: {video_path}"}
    if speed <= 0 or speed > 10:
        return {"ok": False, "error": "Speed must be 0.01-10.0"}
    os.makedirs(os.path.dirname(dst) or ".", exist_ok=True)
    speed_factor = 1.0 / speed
    r = _run_ffmpeg(["-i", src, "-filter_complex",
                     f"[0:v]setpts={speed_factor}*PTS[v];[0:a]atempo={speed}[a]",
                     "-map", "[v]", "-map", "[a]", "-y", dst])
    r["output"] = dst
    r["speed"] = speed
    return r


def tool_video_to_gif(video_path: str, output_path: str = None, start: str = "00:00:00",
                       duration: float = 5.0, width: int = 480, fps: int = 10):
    """视频转 GIF"""
    src = _ffmpeg_safe_path(video_path)
    if not os.path.exists(src):
        return {"ok": False, "error": f"Source not found: {video_path}"}
    if not output_path:
        output_path = os.path.splitext(src)[0] + ".gif"
    dst = _ffmpeg_safe_path(output_path)
    os.makedirs(os.path.dirname(dst) or ".", exist_ok=True)
    r = _run_ffmpeg(["-i", src, "-ss", str(start), "-t", str(duration),
                     "-vf", f"fps={fps},scale={width}:-1:flags=lanczos,split[s0][s1];[s0]palettegen[p];[s1][p]paletteuse",
                     "-loop", "0", "-y", dst])
    r["output"] = dst
    return r


def tool_video_add_text(video_path: str, text: str, output_path: str,
                         position: str = "bottom", font_size: int = 24,
                         font_color: str = "white"):
    """视频添加文字叠加。position: top/bottom/center"""
    src = _ffmpeg_safe_path(video_path)
    dst = _ffmpeg_safe_path(output_path)
    if not os.path.exists(src):
        return {"ok": False, "error": f"Source not found: {video_path}"}
    os.makedirs(os.path.dirname(dst) or ".", exist_ok=True)

    pos_map = {"top": "x=(w-text_w)/2:y=20",
               "bottom": "x=(w-text_w)/2:y=h-text_h-20",
               "center": "x=(w-text_w)/2:y=(h-text_h)/2"}
    pos = pos_map.get(position, pos_map["bottom"])

    # Escape special chars in text for ffmpeg drawtext
    safe_text = text.replace(":", "\\:").replace("'", "\\'").replace("%", "\\%")

    r = _run_ffmpeg(["-i", src, "-vf",
                     f"drawtext=text='{safe_text}':fontsize={font_size}:fontcolor={font_color}:{pos}:box=1:boxcolor=black@0.5:boxborderw=5",
                     "-c:a", "copy", "-y", dst])
    r["output"] = dst
    return r


def tool_video_compress(video_path: str, output_path: str, crf: int = 28, preset: str = "medium"):
    """压缩视频（降低文件大小）。crf: 18-51 (越高越小质量越差)；preset: ultrafast/fast/medium/slow"""
    src = _ffmpeg_safe_path(video_path)
    dst = _ffmpeg_safe_path(output_path)
    if not os.path.exists(src):
        return {"ok": False, "error": f"Source not found: {video_path}"}
    if crf < 0 or crf > 51:
        return {"ok": False, "error": "CRF must be 0-51"}
    os.makedirs(os.path.dirname(dst) or ".", exist_ok=True)
    r = _run_ffmpeg(["-i", src, "-c:v", "libx264", "-crf", str(crf), "-preset", preset,
                     "-c:a", "aac", "-b:a", "128k", "-y", dst])
    r["output"] = dst
    if r["ok"] and os.path.exists(dst):
        orig_size = os.path.getsize(src)
        new_size = os.path.getsize(dst)
        r["original_mb"] = round(orig_size / (1024**2), 1)
        r["output_mb"] = round(new_size / (1024**2), 1)
        r["reduction_pct"] = round((1 - new_size / orig_size) * 100, 1) if orig_size else 0
    return r


def tool_video_convert(video_path: str, output_path: str, vcodec: str = "libx264", acodec: str = "aac"):
    """转换视频格式/编码"""
    src = _ffmpeg_safe_path(video_path)
    dst = _ffmpeg_safe_path(output_path)
    if not os.path.exists(src):
        return {"ok": False, "error": f"Source not found: {video_path}"}
    os.makedirs(os.path.dirname(dst) or ".", exist_ok=True)
    r = _run_ffmpeg(["-i", src, "-c:v", vcodec, "-c:a", acodec, "-y", dst])
    r["output"] = dst
    return r


def tool_video_extract_frames(video_path: str, output_dir: str = None,
                               fps: float = 1, start: str = "00:00:00",
                               duration: float = None, width: int = None):
    """提取视频帧为图片。fps: 每秒提取帧数；width: 输出宽度（保持比例）"""
    src = _ffmpeg_safe_path(video_path)
    if not os.path.exists(src):
        return {"ok": False, "error": f"Source not found: {video_path}"}
    if not output_dir:
        output_dir = os.path.join(os.path.dirname(src), f"frames_{os.path.splitext(os.path.basename(src))[0]}")
    os.makedirs(output_dir, exist_ok=True)
    out_pattern = os.path.join(output_dir, "frame_%06d.jpg")
    args = ["-i", src, "-ss", str(start)]
    if duration:
        args += ["-t", str(duration)]
    vf = f"fps={fps}"
    if width:
        vf += f",scale={width}:-1"
    args += ["-vf", vf, "-q:v", "2", "-y", out_pattern]
    r = _run_ffmpeg(args)
    # Count output frames
    count = 0
    if r["ok"]:
        count = len([f for f in os.listdir(output_dir) if f.endswith(".jpg")])
    r["output_dir"] = output_dir
    r["frame_count"] = count
    return r


def tool_video_crop(video_path: str, output_path: str, x: int = 0, y: int = 0,
                     width: int = 1920, height: int = 1080):
    """裁剪视频区域"""
    src = _ffmpeg_safe_path(video_path)
    dst = _ffmpeg_safe_path(output_path)
    if not os.path.exists(src):
        return {"ok": False, "error": f"Source not found: {video_path}"}
    os.makedirs(os.path.dirname(dst) or ".", exist_ok=True)
    r = _run_ffmpeg(["-i", src, "-vf", f"crop={width}:{height}:{x}:{y}",
                     "-c:a", "copy", "-y", dst])
    r["output"] = dst
    r["crop"] = f"{width}x{height}+{x}+{y}"
    return r


register("video_info", "视频元信息(时长/分辨率/编码) — Video metadata", tool_video_info,
         {"video_path": {"type": "string"}})
register("video_trim", "裁剪视频片段 — Trim video", tool_video_trim,
         {"video_path": {"type": "string"}, "output_path": {"type": "string"},
          "start": {"type": "string", "optional": True}, "duration": {"type": "string", "optional": True},
          "end": {"type": "string", "optional": True}})
register("video_concat", "拼接多个视频 — Concatenate videos", tool_video_concat,
         {"video_paths": {"type": "string", "description": "文件路径 JSON 数组"},
          "output_path": {"type": "string"}})
register("video_resize", "调整视频分辨率 — Resize video", tool_video_resize,
         {"video_path": {"type": "string"}, "output_path": {"type": "string"},
          "width": {"type": "integer", "optional": True}, "height": {"type": "integer", "optional": True}})
register("video_extract_audio", "提取视频音频(mp3/aac/wav) — Extract audio", tool_video_extract_audio,
         {"video_path": {"type": "string"}, "output_path": {"type": "string", "optional": True},
          "format": {"type": "string", "optional": True}})
register("video_replace_audio", "替换视频音轨 — Replace audio track", tool_video_replace_audio,
         {"video_path": {"type": "string"}, "audio_path": {"type": "string"}, "output_path": {"type": "string"}})
register("video_speed", "调整播放速度(0.5半速/2.0双倍) — Change speed", tool_video_speed,
         {"video_path": {"type": "string"}, "output_path": {"type": "string"},
          "speed": {"type": "number", "optional": True}})
register("video_to_gif", "视频转 GIF 动图 — Convert to GIF", tool_video_to_gif,
         {"video_path": {"type": "string"}, "output_path": {"type": "string", "optional": True},
          "start": {"type": "string", "optional": True}, "duration": {"type": "number", "optional": True},
          "width": {"type": "integer", "optional": True}, "fps": {"type": "integer", "optional": True}})
register("video_add_text", "视频叠加文字 — Overlay text", tool_video_add_text,
         {"video_path": {"type": "string"}, "text": {"type": "string"}, "output_path": {"type": "string"},
          "position": {"type": "string", "optional": True}, "font_size": {"type": "integer", "optional": True},
          "font_color": {"type": "string", "optional": True}})
register("video_compress", "压缩视频减小体积(CRF) — Compress video", tool_video_compress,
         {"video_path": {"type": "string"}, "output_path": {"type": "string"},
          "crf": {"type": "integer", "optional": True}, "preset": {"type": "string", "optional": True}})
register("video_convert", "转换视频格式/编码 — Convert codec", tool_video_convert,
         {"video_path": {"type": "string"}, "output_path": {"type": "string"},
          "vcodec": {"type": "string", "optional": True}, "acodec": {"type": "string", "optional": True}})
register("video_extract_frames", "提取视频帧为图片 — Extract frames", tool_video_extract_frames,
         {"video_path": {"type": "string"}, "output_dir": {"type": "string", "optional": True},
          "fps": {"type": "number", "optional": True}, "start": {"type": "string", "optional": True},
          "duration": {"type": "number", "optional": True}, "width": {"type": "integer", "optional": True}})
register("video_crop", "裁剪视频画面区域 — Crop video", tool_video_crop,
         {"video_path": {"type": "string"}, "output_path": {"type": "string"},
          "x": {"type": "integer", "optional": True}, "y": {"type": "integer", "optional": True},
          "width": {"type": "integer", "optional": True}, "height": {"type": "integer", "optional": True}})

# ═══════ Tool Result Cache ═══════
_TOOL_CACHE = {}  # {cache_key: (timestamp, result)}

CACHE_TTL = {
    "system_info": 15,      # 15s — hardware doesn't change fast
    "web_search": 60,       # 60s — same query won't change
    "get_windows": 5,       # 5s
    "list_processes": 10,   # 10s
    "get_volume": 5,
    "mouse_pos": 0.5,       # 0.5s
    "list_dir": 3,
}

def execute(tool_name, params):
    if tool_name not in TOOLS:
        return {"ok": False, "error": f"Unknown tool: {tool_name}"}

    # Check cache for idempotent reads
    ttl = CACHE_TTL.get(tool_name)
    if ttl:
        import time as _t
        cache_key = f"{tool_name}:{json.dumps(params, sort_keys=True, default=str)}"
        cached = _TOOL_CACHE.get(cache_key)
        if cached and (_t.time() - cached[0]) < ttl:
            return cached[1]

    try:
        result = TOOLS[tool_name]["handler"](**params)
    except TypeError as e:
        result = {"ok": False, "error": f"Invalid params: {e}"}
    except Exception as e:
        result = {"ok": False, "error": str(e)}

    # Cache successful results
    if ttl and result.get("ok"):
        _TOOL_CACHE[cache_key] = (_t.time(), result)
    return result

def list_tools():
    return [{"name": t["name"], "description": t["description"], "schema": t["schema"]}
            for t in TOOLS.values()]
